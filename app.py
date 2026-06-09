print("🚀 CLEAN RAG APP") 
import gradio as gr 
import traceback 
import mlflow 
import mlflow_tracker as mlt 
import time 
from langchain_text_splitters import RecursiveCharacterTextSplitter 
from langchain_community.vectorstores import FAISS 
from langchain_core.documents import Document 
from langchain_community.embeddings import FastEmbedEmbeddings
from sentence_transformers import CrossEncoder 
from groq import Groq 
from pptx import Presentation
from dotenv import dotenv_values
import os

config = dotenv_values(".env")

client = Groq(
    api_key=config["GROQ_API_KEY"]
)

print("CLIENT CREATED SUCCESSFULLY")
reranker = CrossEncoder(
    "cross-encoder/ms-marco-MiniLM-L-6-v2"
)

vectorstore = None

PROMPTS = {
    "v1": "Answer ONLY from context.\n\nContext:\n{context}\n\nQuestion:\n{query}",
    "v2": "Strict RAG. If not found say I don't know.\n\nContext:\n{context}\n\nQuery:\n{query}"
}


# -------- DOCUMENT LOADER --------
def load_document(file):
    path = file.name

    if path.endswith(".pdf"):
        from langchain_community.document_loaders import PyPDFLoader
        return PyPDFLoader(path).load()

    elif path.endswith(".pptx"):
        prs = Presentation(path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return [Document(page_content=text)]

    else:
        raise ValueError("Only PDF/PPTX supported")


# -------- PROCESS --------
def process_file(file, chunk_size, chunk_overlap):
    global vectorstore

    try:
        if file is None:
            return "Upload file first"

        #  MLflow START
        mlt.start_run("Processing")

        mlt.log_params({
            "chunk_size": chunk_size,
            "chunk_overlap": chunk_overlap
        })

        docs = load_document(file)

        splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap
        )

        chunks = splitter.split_documents(docs)

        embedding = FastEmbedEmbeddings()
        vectorstore = FAISS.from_texts(
            [c.page_content for c in chunks],
            embedding
        )

        mlt.log_metrics({"num_chunks": len(chunks)})

        mlt.end_run()

        return f"Processed {len(chunks)} chunks"

    except Exception:
        mlt.end_run()
        return traceback.format_exc()


# -------- RETRIEVE --------
@mlflow.trace(span_type="RETRIEVER")
def retrieve(query):
    return vectorstore.similarity_search_with_score(query, k=10)


# -------- RERANK --------
def rerank(query, docs_scores):
    docs = [d for d, _ in docs_scores]

    pairs = [(query, d.page_content) for d in docs]
    scores = reranker.predict(pairs)

    import math
    results = []

    for (doc, sim), rerank_score in zip(docs_scores, scores):
        sim = 1 / (1 + sim)
        rerank_score = 1 / (1 + math.exp(-rerank_score))

        score = 0.7 * rerank_score + 0.3 * sim
        results.append((doc, score))

    return sorted(results, key=lambda x: x[1], reverse=True)[:3]


# -------- GENERATE --------
@mlflow.trace
def generate(prompt):
    try:
        print("\n" + "="*50)
        print("PROMPT LENGTH:", len(prompt))
        print("PROMPT PREVIEW:")
        print(prompt[:500])
        print("="*50)
        print("GENERATE CALLED")
        res = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {"role": "user", "content": prompt}
            ],
            temperature=0,
            max_tokens=1024
        )

        return res.choices[0].message.content

    except Exception as e:
        print("GENERATE ERROR:")
        print(type(e))
        print(str(e))
        raise


# -------- QA --------
@mlflow.trace
def ask(query, prompt_version):
    global vectorstore

    try:
        if vectorstore is None:
            return "Process file first", "", ""

        mlt.start_run("Query")

        start = mlt.start_timer()

        mlt.log_params({
            "query": query,
            "prompt_version": prompt_version
        })

        docs_scores = retrieve(query)
        ranked = rerank(query, docs_scores)

        context = "\n\n".join([d.page_content for d, _ in ranked])

        prompt = PROMPTS[prompt_version].format(
            context=context,
            query=query
        )

        answer = generate(prompt)

        latency = mlt.end_timer(start)

        mlt.log_metrics({
            "latency": latency,
            **{f"confidence_{i}": float(s) for i, (_, s) in enumerate(ranked)}
        })

        mlt.log_text(answer)

        mlt.end_run()

        sources = "\n\n".join([d.page_content[:200] for d, _ in ranked])
        conf = "\n".join([str(round(s, 3)) for _, s in ranked])

        return answer, sources, conf

    except Exception:
        mlt.end_run()
        return traceback.format_exc(), "", ""


# -------- UI --------
with gr.Blocks() as demo:
    gr.Markdown("# RAG + MLflow (Clean Architecture)")

    with gr.Row():

        with gr.Column():
            file = gr.File()
            cs = gr.Slider(300, 1500, value=800)
            co = gr.Slider(0, 300, value=100)
            btn = gr.Button("Process")
            status = gr.Textbox()

        with gr.Column():
            q = gr.Textbox()
            pv = gr.Dropdown(["v1", "v2"], value="v1")

            ans = gr.Textbox(lines=6)
            src = gr.Textbox(lines=8)
            conf = gr.Textbox()

            ask_btn = gr.Button("Ask")

    btn.click(process_file, [file, cs, co], status)
    ask_btn.click(ask, [q, pv], [ans, src, conf])

demo.launch()